"""
builder.py
==========
Erstellt die PowerPoint-Praesentation aus den geparsten DataFrames.
"""

import sys

import pandas as pd

from .converter import add_emu_columns, prepare_dataframe


def _add_text_in_form(slide, row, MSO_SHAPE, Emu, Pt, RGBColor, PP_ALIGN) -> bool:
    """Fuegt ein Element 'Text in Form' (z. B. beschriftetes Rechteck) hinzu."""
    try:
        x      = Emu(int(row.get("x_emu",      0)))
        y      = Emu(int(row.get("y_emu",      0)))
        width  = Emu(int(row.get("breite_emu", 500000)))
        height = Emu(int(row.get("hoehe_emu",  200000)))

        form   = str(row.get("Form_Typ", "Rechteck")).lower()
        stype  = MSO_SHAPE.OVAL if "kreis" in form else MSO_SHAPE.RECTANGLE
        rgb_f  = row.get("rgb_form", (200, 200, 200))
        rgb_t  = row.get("rgb_text",  (0, 0, 0))
        text   = str(row.get("Textinhalt", ""))

        shape  = slide.shapes.add_shape(stype, x, y, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*rgb_f)
        shape.line.color.rgb      = RGBColor(*rgb_f)

        tf = shape.text_frame
        tf.word_wrap = True
        p  = tf.paragraphs[0]
        p.text           = text
        p.font.size      = Pt(12)
        p.font.color.rgb = RGBColor(*rgb_t)
        p.alignment      = PP_ALIGN.CENTER
        return True
    except Exception as e:
        print(f"    [ERR] Text-in-Form: {e}")
        return False


def _add_textbox(slide, row, Emu, Pt, RGBColor) -> bool:
    """Fuegt ein reines Textfeld ohne Form-Hintergrund hinzu."""
    try:
        x      = Emu(int(row.get("x_emu",      0)))
        y      = Emu(int(row.get("y_emu",      0)))
        width  = Emu(int(row.get("breite_emu", 500000)))
        height = Emu(int(row.get("hoehe_emu",  200000)))
        rgb_t  = row.get("rgb_text", (0, 0, 0))
        text   = str(row.get("Textinhalt", ""))

        tb = slide.shapes.add_textbox(x, y, width, height)
        tf = tb.text_frame
        tf.word_wrap     = True
        p  = tf.paragraphs[0]
        p.text           = text
        p.font.size      = Pt(12)
        p.font.color.rgb = RGBColor(*rgb_t)
        return True
    except Exception as e:
        print(f"    [ERR] Textbox: {e}")
        return False


def _add_shape_only(slide, row, MSO_SHAPE, Emu, RGBColor) -> bool:
    """Fuegt eine Form ohne Text hinzu."""
    try:
        x      = Emu(int(row.get("x_emu",      0)))
        y      = Emu(int(row.get("y_emu",      0)))
        width  = Emu(int(row.get("breite_emu", 300000)))
        height = Emu(int(row.get("hoehe_emu",  300000)))

        form   = str(row.get("Form_Typ", "Rechteck")).lower()
        stype  = MSO_SHAPE.OVAL if "kreis" in form else MSO_SHAPE.RECTANGLE
        rgb_f  = row.get("rgb_form", (150, 150, 150))

        shape  = slide.shapes.add_shape(stype, x, y, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*rgb_f)
        shape.line.color.rgb      = RGBColor(*rgb_f)
        return True
    except Exception as e:
        print(f"    [ERR] Shape: {e}")
        return False


def create_powerpoint(dataframes: dict, output_path: str) -> str:
    """
    Erstellt eine PowerPoint-Datei aus den geparsten DataFrames.

    Parameters
    ----------
    dataframes  : dict  – Ausgabe von ``parse_llm_output()``
    output_path : str   – Zieldatei (z. B. "output.pptx")

    Returns
    -------
    str – Absoluter Pfad zur gespeicherten Datei
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Emu, Inches, Pt
    except ImportError:
        sys.exit("Installiere: pip install python-pptx")

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Leere Folie

    success, errors = 0, 0

    # Reihenfolge: C (Text in Form) -> B (Reiner Text) -> D (Formen ohne Text)
    tasks = [
        ("TEIL C", "Text in Form",       lambda r: _add_text_in_form(slide, r, MSO_SHAPE, Emu, Pt, RGBColor, PP_ALIGN)),
        ("TEIL B", "Reiner Text",        lambda r: _add_textbox(slide, r, Emu, Pt, RGBColor)),
        ("TEIL D", "Formen (kein Text)", lambda r: _add_shape_only(slide, r, MSO_SHAPE, Emu, RGBColor)),
    ]

    for teil, label, fn in tasks:
        df = dataframes.get(teil)
        if df is None:
            continue
        df = add_emu_columns(df)
        df = prepare_dataframe(df)
        print(f"\n[{label}]: {len(df)} Elemente ...")
        for _, row in df.iterrows():
            eid  = row.get("Element_ID", "?")
            text = str(row.get("Textinhalt", ""))[:40]
            if fn(row):
                print(f"  [OK] {eid}: {text}")
                success += 1
            else:
                errors += 1

    prs.save(output_path)
    print(f"\n{'='*58}")
    print(f"PowerPoint gespeichert : {output_path}")
    print(f"Erfolgreiche Elemente  : {success}")
    print(f"Fehlerhafte Elemente   : {errors}")
    print(f"{'='*58}")
    return output_path
